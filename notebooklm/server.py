import json
import re
import threading
import base64
import math
import os
import errno
import socket
from io import BytesIO
from collections import Counter

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
from http import HTTPStatus
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from urllib.parse import urlparse, quote_plus

BASE_DIR = Path(__file__).resolve().parent
DATA_FILE = BASE_DIR / "data.json"
LOCK = threading.Lock()
STOPWORDS = {
    "the", "and", "for", "that", "with", "from", "this", "have", "are", "was", "were", "been", "into", "about",
    "your", "their", "they", "them", "will", "would", "could", "should", "there", "what", "when", "where", "which",
    "while", "also", "than", "then", "after", "before", "between", "through", "using", "used", "use", "into", "over",
    "under", "most", "more", "some", "many", "such", "each", "other", "only", "very", "much", "like", "just",
}
CHUNK_SIZE_TOKENS = 90
CHUNK_OVERLAP_TOKENS = 20


def default_state():
    return {"notebooks": []}


def load_state():
    with LOCK:
        if not DATA_FILE.exists():
            state = default_state()
            DATA_FILE.write_text(json.dumps(state, indent=2), encoding="utf-8")
            return state
        try:
            payload = json.loads(DATA_FILE.read_text(encoding="utf-8"))
            if not isinstance(payload, dict):
                return default_state()
            notebooks = payload.get("notebooks", [])
            if not isinstance(notebooks, list):
                notebooks = []
            return {"notebooks": notebooks}
        except json.JSONDecodeError:
            state = default_state()
            DATA_FILE.write_text(json.dumps(state, indent=2), encoding="utf-8")
            return state


def save_state(state):
    with LOCK:
        DATA_FILE.write_text(json.dumps(state, indent=2), encoding="utf-8")


def safe_text(value):
    return str(value or "")


def split_sentences(text):
    normalized = safe_text(text).replace("\r", "\n")
    out = []
    for line in normalized.split("\n"):
        for part in re.split(r"[.!?]+", line):
            sentence = part.strip()
            if sentence:
                out.append(sentence)
    return out


def tokenize(text):
    return re.findall(r"[a-zA-Z0-9]+", safe_text(text).lower())


def has_readable_text(text):
    tokens = tokenize(text)
    return len(tokens) >= 3


def content_tokens(text):
    return [t for t in tokenize(text) if len(t) > 2 and t not in STOPWORDS]


def extract_base64_text(b64_data, kind):
    if not b64_data:
        return ""
    try:
        b64_data = b64_data.strip()
        padding = len(b64_data) % 4
        if padding:
            b64_data += '=' * (4 - padding)
        
        file_data = base64.b64decode(b64_data)
        if kind == "pdf" and PdfReader:
            reader = PdfReader(BytesIO(file_data))
            return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        elif kind == "docx" and Document:
            doc = Document(BytesIO(file_data))
            return "\n".join(para.text for para in doc.paragraphs)
        elif kind == "pptx" and Presentation:
            prs = Presentation(BytesIO(file_data))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            return "\n".join(text)
    except Exception as e:
        print(f"Extraction error for {kind}: {e}")
        return ""
    return ""


def normalize_sources(sources):
    normalized = []
    for source in sources or []:
        if not isinstance(source, dict):
            continue
        
        kind = safe_text(source.get("kind") or "file")
        text = safe_text(source.get("text"))
        b64 = safe_text(source.get("base64"))
        
        if not text.strip() and b64:
            extracted = extract_base64_text(b64, kind)
            if extracted:
                text = extracted

        normalized.append(
            {
                "name": safe_text(source.get("name") or "source"),
                "kind": kind,
                "text": text,
            }
        )
    return normalized


def chunk_source_text(text, max_tokens=CHUNK_SIZE_TOKENS, overlap_tokens=CHUNK_OVERLAP_TOKENS):
    chunks = []
    current = []
    current_token_count = 0

    for sentence in split_sentences(text):
        sentence = sentence.strip()
        sentence_tokens = tokenize(sentence)
        if len(sentence_tokens) < 3:
            continue

        if current and current_token_count + len(sentence_tokens) > max_tokens:
            chunk = ". ".join(current).strip()
            if chunk:
                if chunk[-1] not in ".!?":
                    chunk += "."
                chunks.append(chunk)

            if overlap_tokens > 0:
                overlap = []
                overlap_count = 0
                for prev in reversed(current):
                    overlap.insert(0, prev)
                    overlap_count += len(tokenize(prev))
                    if overlap_count >= overlap_tokens:
                        break
                current = overlap
                current_token_count = sum(len(tokenize(x)) for x in current)
            else:
                current = []
                current_token_count = 0

        current.append(sentence)
        current_token_count += len(sentence_tokens)

    if current:
        chunk = ". ".join(current).strip()
        if chunk:
            if chunk[-1] not in ".!?":
                chunk += "."
            chunks.append(chunk)

    if not chunks and has_readable_text(text):
        fallback = safe_text(text).strip()
        if fallback:
            chunks.append(fallback[:1200])
    return chunks


def build_retrieval_items(sources):
    items = []
    for source in normalize_sources(sources):
        for chunk in chunk_source_text(source["text"]):
            tokens = content_tokens(chunk)
            if not tokens:
                continue
            items.append(
                {
                    "source": source["name"],
                    "sentence": chunk,
                    "tokens": tokens,
                }
            )
    return items


def build_idf(items):
    doc_count = len(items)
    frequencies = Counter()
    for item in items:
        frequencies.update(set(item["tokens"]))
    return {term: math.log((doc_count + 1) / (count + 1)) + 1.0 for term, count in frequencies.items()}


def rank_sentences(question, sources, top_n=6):
    items = build_retrieval_items(sources)
    if not items:
        return []

    q_tokens = content_tokens(question)
    if not q_tokens:
        return [
            {"score": 1, "source": item["source"], "sentence": item["sentence"]}
            for item in items[:top_n]
        ]

    idf = build_idf(items)
    q_counts = Counter(q_tokens)
    q_total = max(1, len(q_tokens))
    q_norm = 0.0
    for term, count in q_counts.items():
        weight = (count / q_total) * idf.get(term, 1.0)
        q_norm += weight * weight
    q_norm = math.sqrt(q_norm)

    scored = []
    for item in items:
        d_counts = Counter(item["tokens"])
        d_total = max(1, len(item["tokens"]))

        dot = 0.0
        d_norm_sq = 0.0
        for term, count in d_counts.items():
            d_w = (count / d_total) * idf.get(term, 1.0)
            d_norm_sq += d_w * d_w

        if q_norm > 0 and d_norm_sq > 0:
            for term, q_count in q_counts.items():
                d_count = d_counts.get(term, 0)
                if not d_count:
                    continue
                weight = idf.get(term, 1.0)
                q_w = (q_count / q_total) * weight
                d_w = (d_count / d_total) * weight
                dot += q_w * d_w
            cosine = dot / (q_norm * math.sqrt(d_norm_sq))
        else:
            cosine = 0.0

        overlap = len(set(q_counts.keys()).intersection(d_counts.keys()))
        overlap_boost = overlap / max(1, len(q_counts))
        score = cosine + (0.15 * overlap_boost)

        if score > 0:
            scored.append({"score": score, "source": item["source"], "sentence": item["sentence"]})

    if not scored:
        return [
            {"score": 1, "source": item["source"], "sentence": item["sentence"]}
            for item in items[:top_n]
        ]

    scored.sort(key=lambda x: x["score"], reverse=True)
    return scored[:top_n]


def unique_citations(items):
    cites = []
    for item in items:
        src = safe_text(item.get("source"))
        if src and src not in cites:
            cites.append(src)
    return cites


def concise_line(text, max_len=220):
    text = safe_text(text).strip()
    if not text:
        return ""
    parts = split_sentences(text)
    line = parts[0].strip() if parts else text
    line = re.sub(r"\s+", " ", line)
    if len(line) > max_len:
        line = line[: max_len - 1].rstrip() + "…"
    return line


def format_structured_answer(question, items):
    lines = ["Structured Answer", ""]
    q = safe_text(question).strip()
    if q:
        lines.extend(["Question", f"- {q}", ""])

    lines.append("Key Points")
    source_map = {}
    for idx, item in enumerate(items, start=1):
        point = concise_line(item.get("sentence"))
        if not point:
            continue
        lines.append(f"{idx}. {point}")
        src = safe_text(item.get("source"))
        if src:
            source_map.setdefault(src, []).append(str(idx))

    if source_map:
        lines.extend(["", "Source Mapping"])
        for src, refs in source_map.items():
            lines.append(f"- {src}: points {', '.join(refs)}")
    return "\n".join(lines)


def extraction_diagnostics(sources):
    normalized = normalize_sources(sources)
    if not normalized:
        return "No sources were provided."

    has_binary = False
    no_text_sources = []
    for src in normalized:
        if has_readable_text(src.get("text")):
            continue
        no_text_sources.append(f"{src.get('name', 'source')} ({safe_text(src.get('kind')).upper() or 'FILE'})")
        if safe_text(src.get("kind")).lower() in {"pdf", "docx", "pptx"}:
            has_binary = True

    if not no_text_sources:
        return "Readable text exists, but retrieval found no strong matches."

    reason = "No readable text could be extracted from one or more uploaded sources."
    if has_binary:
        reason += " Some files appear to be image-based or text extraction returned empty content."
    sample = ", ".join(no_text_sources[:3])
    more = "" if len(no_text_sources) <= 3 else f" (+{len(no_text_sources) - 3} more)"
    return f"{reason} Affected: {sample}{more}."


def normalize_url(candidate):
    url = safe_text(candidate).strip()
    if not url:
        return ""
    if not re.match(r"^[a-zA-Z][a-zA-Z0-9+.-]*://", url):
        url = "https://" + url
    parsed = urlparse(url)
    if not parsed.netloc:
        return ""
    return url


def detect_tool_call(question):
    text = safe_text(question).strip()
    if not text:
        return None

    lower = text.lower()
    if lower.startswith("/youtube "):
        query = text[9:].strip()
        return {"tool": "youtube", "query": query} if query else None
    if lower.startswith("/browser "):
        query = text[9:].strip()
        return {"tool": "browser", "query": query} if query else None
    if lower.startswith("/search "):
        query = text[8:].strip()
        return {"tool": "browser", "query": query} if query else None
    if lower.startswith("/open "):
        url = text[6:].strip()
        return {"tool": "open", "url": url} if url else None

    m = re.match(r"^(?:search|find|look up)\s+(?:on\s+)?youtube\s+(?:for\s+)?(.+)$", text, flags=re.IGNORECASE)
    if m:
        return {"tool": "youtube", "query": m.group(1).strip()}

    m = re.match(r"^(?:search|browse|google|look up|find)\s+(?:for\s+)?(.+)$", text, flags=re.IGNORECASE)
    if m:
        return {"tool": "browser", "query": m.group(1).strip()}

    m = re.match(r"^open\s+(.+)$", text, flags=re.IGNORECASE)
    if m:
        return {"tool": "open", "url": m.group(1).strip()}
    return None


def execute_tool_call(call):
    if not isinstance(call, dict):
        return None
    tool = safe_text(call.get("tool")).lower()

    if tool == "youtube":
        query = safe_text(call.get("query")).strip()
        if not query:
            return None
        url = f"https://www.youtube.com/results?search_query={quote_plus(query)}"
        return {
            "text": f"Tool call: YouTube search for \"{query}\".\nOpen this link: {url}",
            "citations": ["Tool: YouTube"],
            "tool_call": {"name": "youtube", "target_url": url},
        }

    if tool == "browser":
        query = safe_text(call.get("query")).strip()
        if not query:
            return None
        url = f"https://www.google.com/search?q={quote_plus(query)}"
        return {
            "text": f"Tool call: Browser search for \"{query}\".\nOpen this link: {url}",
            "citations": ["Tool: Browser"],
            "tool_call": {"name": "browser", "target_url": url},
        }

    if tool == "open":
        url = normalize_url(call.get("url"))
        if not url:
            return {
                "text": "Tool call failed: invalid URL. Use a valid domain or URL like `open example.com`.",
                "citations": ["Tool: Open URL"],
            }
        return {
            "text": f"Tool call: Open URL.\nOpen this link: {url}",
            "citations": ["Tool: Open URL"],
            "tool_call": {"name": "open_url", "target_url": url},
        }
    return None


def answer_question(question, sources):
    tool_response = execute_tool_call(detect_tool_call(question))
    if tool_response:
        return tool_response

    if re.search(r"\b(mcq|quiz|multiple\s+choice)\b", safe_text(question), flags=re.IGNORECASE):
        if not sources:
            return {
                "text": "No extractable text found. Upload text-based files (txt, md, csv, json, html, xml) or PDF, PPTX or DOCX documents to extract text.",
                "citations": [],
            }
        text, citations = build_quiz(sources)
        return {"text": text, "citations": citations}

    if not sources:
        return {
            "text": "No extractable text found. Upload text-based files (txt, md, csv, json, html, xml) or PDF, PPTX or DOCX documents to extract text.",
            "citations": [],
        }

    top = rank_sentences(question, sources, top_n=5)
    if not top:
        if any(has_readable_text(src.get("text")) for src in normalize_sources(sources)):
            return {
                "text": "I found no direct match in extracted text. Try a narrower question or upload additional sources.",
                "citations": [],
            }
        return {
            "text": extraction_diagnostics(sources) + " Try uploading a text-based TXT/MD/CSV/JSON file or a searchable PDF.",
            "citations": [],
        }

    return {
        "text": format_structured_answer(question, top),
        "citations": unique_citations(top),
    }


def top_terms(sources, limit=12):
    counts = Counter()
    for source in normalize_sources(sources):
        for token in tokenize(source["text"]):
            if len(token) < 4 or token in STOPWORDS:
                continue
            counts[token] += 1
    return [w for w, _ in counts.most_common(limit)]


def extract_timeline_items(sources, limit=8):
    pattern = re.compile(r"\b(?:\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|\d{4}|(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\s+\d{1,2}(?:,\s*\d{4})?)\b", re.IGNORECASE)
    items = []

    for source in normalize_sources(sources):
        for sentence in split_sentences(source["text"]):
            if len(sentence) < 20:
                continue
            if pattern.search(sentence):
                items.append({"source": source["name"], "sentence": sentence.strip()})
            if len(items) >= limit:
                return items
    return items


def build_briefing(sources):
    key_sentences = rank_sentences("summary key findings impact risks actions", sources, top_n=8)
    if not key_sentences:
        return "No meaningful text found in the uploaded sources.", []

    lines = ["Executive Briefing", "", "Key Findings"]
    lines.extend([f"- {x['sentence']}" for x in key_sentences[:4]])
    lines.append("")
    lines.append("Risks and Unknowns")
    lines.extend([f"- {x['sentence']}" for x in key_sentences[4:6]])
    lines.append("")
    lines.append("Suggested Next Steps")
    lines.extend([f"- Validate: {x['sentence']}" for x in key_sentences[6:8]])

    return "\n".join(lines), unique_citations(key_sentences)


def build_faq(sources):
    key_sentences = rank_sentences("definition purpose methodology metrics constraints results", sources, top_n=10)
    if not key_sentences:
        return "No meaningful text found in the uploaded sources.", []

    questions = [
        "What is the core objective?",
        "What evidence supports it?",
        "What are the major constraints?",
        "What are the main risks?",
        "What should be done next?",
    ]

    lines = ["FAQ"]
    citations = []
    for i, question in enumerate(questions):
        sentence = key_sentences[i % len(key_sentences)]
        lines.append("")
        lines.append(f"Q{i + 1}. {question}")
        lines.append(f"A. {sentence['sentence']}")
        if sentence["source"] not in citations:
            citations.append(sentence["source"])

    return "\n".join(lines), citations


def build_timeline(sources):
    events = extract_timeline_items(sources, limit=10)
    if not events:
        fallback = rank_sentences("milestones timeline sequence", sources, top_n=6)
        if not fallback:
            return "No timeline-like events found in sources.", []
        lines = ["Timeline (inferred order)"]
        lines.extend([f"- {x['sentence']}" for x in fallback])
        return "\n".join(lines), unique_citations(fallback)

    lines = ["Timeline"]
    for idx, ev in enumerate(events, start=1):
        lines.append(f"{idx}. {ev['sentence']}")
    return "\n".join(lines), unique_citations(events)


def build_study_guide(sources):
    terms = top_terms(sources, limit=10)
    key_sentences = rank_sentences("important concepts examples implications", sources, top_n=8)
    if not key_sentences:
        return "No meaningful text found in the uploaded sources.", []

    lines = ["Study Guide", "", "Core Concepts"]
    if terms:
        lines.extend([f"- {term}" for term in terms[:8]])

    lines.append("")
    lines.append("Key Explanations")
    lines.extend([f"- {x['sentence']}" for x in key_sentences[:5]])

    lines.append("")
    lines.append("Self-Test Questions")
    lines.extend(
        [
            "1. Which concept appears most frequently across sources?",
            "2. What assumptions drive the main argument?",
            "3. Which evidence is strongest and why?",
            "4. What would invalidate the current conclusion?",
        ]
    )

    return "\n".join(lines), unique_citations(key_sentences)


def build_flash_cards(sources):
    prompts = rank_sentences("definitions key terms concepts metrics process outcomes", sources, top_n=12)
    if not prompts:
        return "No meaningful text found in the uploaded sources.", []

    lines = ["Flash Cards"]
    citations = []
    card_no = 1
    for item in prompts:
        sentence = safe_text(item.get("sentence")).strip()
        parts = re.split(r"\s*[;:,-]\s*", sentence, maxsplit=1)
        if len(parts) == 2 and len(parts[0].split()) <= 8:
            front = parts[0].strip().rstrip(".")
            back = parts[1].strip()
        else:
            words = sentence.split()
            if len(words) < 6:
                continue
            split_at = min(6, max(3, len(words) // 3))
            front = " ".join(words[:split_at]).strip().rstrip(".")
            back = " ".join(words[split_at:]).strip()

        if not front or not back:
            continue

        lines.append("")
        lines.append(f"{card_no}. Q: {front}?")
        lines.append(f"   A: {back}")
        card_no += 1

        source_name = safe_text(item.get("source"))
        if source_name and source_name not in citations:
            citations.append(source_name)

        if card_no > 10:
            break

    if card_no == 1:
        return "Could not form flash cards from current source text.", []

    return "\n".join(lines), citations


def sentence_to_qa(sentence):
    sentence = safe_text(sentence).strip()
    if not sentence:
        return "", ""
    parts = re.split(r"\s*[;:,-]\s*", sentence, maxsplit=1)
    if len(parts) == 2 and len(parts[0].split()) <= 10 and len(parts[1].split()) >= 3:
        question = f"What is {parts[0].strip().rstrip('.')}?"
        answer = parts[1].strip()
        return question, answer

    words = sentence.split()
    if len(words) < 8:
        return "", ""
    pivot = min(8, max(4, len(words) // 3))
    question = "Explain: " + " ".join(words[:pivot]).strip().rstrip(".") + "?"
    answer = " ".join(words[pivot:]).strip()
    return question, answer


def build_quiz(sources):
    facts = rank_sentences("core concepts definitions process outcomes evidence metrics", sources, top_n=20)
    if not facts:
        return "No meaningful text found in the uploaded sources.", []

    lines = ["Quiz (MCQ)", "", "Choose one correct option (A-D) for each question."]
    key = ["", "Answer Key"]
    citations = []
    all_sentences = []
    for x in facts:
        raw = safe_text(x.get("sentence")).strip()
        if not raw:
            continue
        first = split_sentences(raw)
        candidate = first[0].strip() if first else raw
        if candidate:
            all_sentences.append(candidate)
    concept_bank = top_terms(sources, limit=40)
    qno = 1
    for idx, item in enumerate(facts):
        sentence = safe_text(item.get("sentence")).strip()
        if not sentence:
            continue
        first_parts = split_sentences(sentence)
        sentence = first_parts[0].strip() if first_parts else sentence

        topic_tokens = [t for t in content_tokens(sentence) if len(t) > 3]
        topic = " ".join(topic_tokens[:3]) if topic_tokens else "the source material"
        lines.append(f"{qno}. Which statement is supported by the sources about {topic}?")

        distractors = []
        for other in all_sentences:
            if other == sentence:
                continue
            if other not in distractors:
                distractors.append(other)
            if len(distractors) >= 3:
                break

        while len(distractors) < 3:
            seed = concept_bank[(idx + len(distractors)) % len(concept_bank)] if concept_bank else "the topic"
            distractors.append(f"It is mainly unrelated to {seed}.")

        correct_position = idx % 4
        labels = ["A", "B", "C", "D"]
        options = []
        d_i = 0
        for pos in range(4):
            if pos == correct_position:
                options.append(sentence)
            else:
                options.append(distractors[d_i])
                d_i += 1
        for opt_i, option_text in enumerate(options):
            lines.append(f"   {labels[opt_i]}. {option_text}")
        key.append(f"{qno}. {labels[correct_position]}")

        source_name = safe_text(item.get("source"))
        if source_name and source_name not in citations:
            citations.append(source_name)
        qno += 1
        if qno > 7:
            break

    if qno == 1:
        return "Could not build MCQ quiz questions from current source text.", []

    lines.extend(key)
    return "\n".join(lines), citations


def build_qa_sheet(sources):
    facts = rank_sentences("important points who what why how implications results", sources, top_n=12)
    if not facts:
        return "No meaningful text found in the uploaded sources.", []

    lines = ["Q&A Sheet"]
    citations = []
    qno = 1
    for item in facts:
        question, answer = sentence_to_qa(item.get("sentence"))
        if not question or not answer:
            continue
        lines.append("")
        lines.append(f"Q{qno}. {question}")
        lines.append(f"A{qno}. {answer}")
        source_name = safe_text(item.get("source"))
        if source_name and source_name not in citations:
            citations.append(source_name)
        qno += 1
        if qno > 10:
            break

    if qno == 1:
        return "Could not create Q&A pairs from current source text.", []

    return "\n".join(lines), citations


def generate_studio(task, sources):
    task = safe_text(task).strip().lower()
    if task == "briefing":
        text, citations = build_briefing(sources)
        return {"title": "Briefing Document", "text": text, "citations": citations}
    if task == "faq":
        text, citations = build_faq(sources)
        return {"title": "FAQ", "text": text, "citations": citations}
    if task == "timeline":
        text, citations = build_timeline(sources)
        return {"title": "Timeline", "text": text, "citations": citations}
    if task == "study_guide":
        text, citations = build_study_guide(sources)
        return {"title": "Study Guide", "text": text, "citations": citations}
    if task == "flash_cards":
        text, citations = build_flash_cards(sources)
        return {"title": "Flash Cards", "text": text, "citations": citations}
    if task == "quiz":
        text, citations = build_quiz(sources)
        return {"title": "Quiz", "text": text, "citations": citations}
    if task == "qa_sheet":
        text, citations = build_qa_sheet(sources)
        return {"title": "Q&A Sheet", "text": text, "citations": citations}

    return {
        "title": "Studio output",
        "text": "Unsupported studio task. Use briefing, faq, timeline, study_guide, flash_cards, quiz, or qa_sheet.",
        "citations": [],
    }


class Handler(SimpleHTTPRequestHandler):
    def _json_response(self, payload, status=HTTPStatus.OK):
        body = json.dumps(payload).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _read_json_body(self):
        try:
            length = int(self.headers.get("Content-Length", "0"))
        except ValueError:
            return None
        raw = self.rfile.read(length) if length > 0 else b""
        if not raw:
            return {}
        try:
            return json.loads(raw.decode("utf-8"))
        except json.JSONDecodeError:
            return None

    def do_GET(self):
        parsed = urlparse(self.path)
        if parsed.path == "/api/state":
            self._json_response(load_state())
            return

        if parsed.path == "/":
            self.path = "/index.html"

        return super().do_GET()

    def do_PUT(self):
        parsed = urlparse(self.path)
        if parsed.path != "/api/state":
            self.send_error(HTTPStatus.NOT_FOUND, "Not found")
            return

        payload = self._read_json_body()
        if payload is None:
            self.send_error(HTTPStatus.BAD_REQUEST, "Invalid JSON")
            return

        notebooks = payload.get("notebooks", []) if isinstance(payload, dict) else []
        if not isinstance(notebooks, list):
            self.send_error(HTTPStatus.BAD_REQUEST, "'notebooks' must be a list")
            return

        save_state({"notebooks": notebooks})
        self._json_response({"ok": True})

    def do_POST(self):
        parsed = urlparse(self.path)
        payload = self._read_json_body()
        if payload is None:
            self.send_error(HTTPStatus.BAD_REQUEST, "Invalid JSON")
            return

        if parsed.path == "/api/answer":
            question = safe_text(payload.get("question")) if isinstance(payload, dict) else ""
            sources = payload.get("sources") if isinstance(payload, dict) else []
            if not isinstance(sources, list):
                self.send_error(HTTPStatus.BAD_REQUEST, "'sources' must be a list")
                return
            self._json_response(answer_question(question, sources))
            return

        if parsed.path == "/api/studio":
            task = safe_text(payload.get("task")) if isinstance(payload, dict) else ""
            sources = payload.get("sources") if isinstance(payload, dict) else []
            if not isinstance(sources, list):
                self.send_error(HTTPStatus.BAD_REQUEST, "'sources' must be a list")
                return
            self._json_response(generate_studio(task, sources))
            return

        self.send_error(HTTPStatus.NOT_FOUND, "Not found")


def run(host="0.0.0.0", port=3000):
    def lan_ip():
        # Resolve the active outbound interface IP without sending data.
        sock = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        try:
            sock.connect(("8.8.8.8", 80))
            return sock.getsockname()[0]
        except OSError:
            return "127.0.0.1"
        finally:
            sock.close()

    def normalize_url(value):
        raw = safe_text(value).strip()
        if not raw:
            return ""
        if raw.startswith("http://") or raw.startswith("https://"):
            return raw
        return f"https://{raw}"

    def infer_cloud_url(port_number):
        # Common env vars for cloud IDE public URLs.
        direct = normalize_url(os.environ.get("PUBLIC_URL") or os.environ.get("CODESANDBOX_HOST"))
        if direct:
            return direct

        gitpod_url = normalize_url(os.environ.get("GITPOD_WORKSPACE_URL"))
        if gitpod_url:
            return gitpod_url

        csb_host = safe_text(os.environ.get("CODESANDBOX_HOST")).strip()
        if csb_host:
            host_only = csb_host.split("://", 1)[-1]
            if not host_only.startswith("http"):
                if f"-{port_number}." in host_only:
                    return f"https://{host_only}"
                return f"https://{host_only}"

        return ""

    preferred = int(port)
    candidates = [preferred]
    if preferred == 3000:
        candidates.extend([8000, 8080, 8888])

    for candidate in candidates:
        try:
            server = ThreadingHTTPServer((host, candidate), Handler)
            local_url = f"http://127.0.0.1:{candidate}"
            lan_url = f"http://{lan_ip()}:{candidate}"
            cloud_url = infer_cloud_url(candidate)
            print("Serving Smart Notebook AI")
            print(f"Local:    {local_url}")
            print(f"LAN:      {lan_url}")
            if cloud_url:
                print(f"External: {cloud_url}")
            else:
                print(f"External: {lan_url}")
            server.serve_forever()
            return
        except OSError as exc:
            if exc.errno == errno.EADDRINUSE and candidate != candidates[-1]:
                continue
            raise


if __name__ == "__main__":
    host = os.environ.get("HOST", "0.0.0.0")
    try:
        port = int(os.environ.get("PORT", "3000"))
    except ValueError:
        port = 3000
    run(host=host, port=port)
